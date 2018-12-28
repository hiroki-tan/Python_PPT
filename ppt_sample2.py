# -*- coding: utf-8 -*-
import os
from pptx import Presentation

SLD_LAYOUT_TITLE_AND_CONTENT = 0

# Presentaitonインスタンスの作成
# デフォルトではC:\Python27\Lib\site-packages\pptx\templates\default.pptxを読み取る
prs = Presentation()

# レイアウトの決定
slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]

# スライド作成(Slide)
slide = prs.slides.add_slide(slide_layout)
## テキストの設定(placeholders)
slide.shapes.title.text = "placeholders[0]"
slide.placeholders[1].text = "placeholders[1]"
print type(slide)
print "len(slide.placeholders) : " + str(len(slide.placeholders))
print slide.placeholders[0].text
print slide.placeholders[1].text
print

print type(slide.shapes.placeholders[1].text_frame)
print slide.shapes.placeholders[0].text_frame.text
print slide.shapes.placeholders[1].text_frame.text

### プレースホルダ内にテキストを追加(TextFrame)
text_frame = slide.shapes.placeholders[1].text_frame
paragraph = text_frame.add_paragraph()
paragraph.text = "add_paragraph text"
#### 段落のレベル
paragraph.level = 1

paragraph = slide.shapes.placeholders[1].text_frame.add_paragraph()
paragraph.text = "add_paragraph text2"
paragraph.level = 8

text_frame.fit_text(font_family='Calibri', max_size=33, bold=True, italic=True, font_file=None)

### プレースホルダを削除
#text_frame.clear()

# スライド追加
slide = prs.slides.add_slide(slide_layout)
## テキストの設定
slide.shapes.title.text = "placeholders[0]"
slide.placeholders[1].text = "placeholders[1]"

# スライドの保存
prs.save('python.pptx')
